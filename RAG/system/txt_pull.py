
from io import StringIO
from pdfminer.high_level import extract_text_to_fp
import pandas as pd
from pptx import Presentation
import subprocess
from docx import Document
from bot import ai_assist, upload, cleanup

def not_implemented(file_path):
    return "File type not implemented" + file_path

def pdf_to_text(file_path):
    print("PDF file detected, processing...")

    output_string = StringIO()
    with open(file_path, 'rb') as f:
        extract_text_to_fp(f, output_string)
    return output_string.getvalue()

def docx_to_text(file_path):
    print("docx file detected, processing...")

    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def txt_to_text(file_path):
    print("txt file detected, processing...")
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def ppt_to_text(file_path):
    print("PPT file detected, processing...")

    p = subprocess.run(['catppt', file_path],
        stdout=subprocess.PIPE,
        stderr=subprocess.DEVNULL)
    return p.stdout.decode('utf-8', errors='ignore')

def doc_to_text(file_path):
    print("doc file detected, processing...")

    p = subprocess.run(['antiword', file_path],
        stdout=subprocess.PIPE,
        stderr=subprocess.DEVNULL)
    return p.stdout.decode('utf-8', errors='ignore')

def csv_to_text(file_path):
    print("CSV file detected, processing...")

    sep = ","
    df = pd.read_csv(file_path, sep=sep, encoding='utf-8', dtype=str, keep_default_na=False)
    lines = []
    for row in df.values:
        lines.append(sep.join(row))
    return "\n".join(lines)

def pptx_to_text(file_path):
    print("PPTX file detected, processing...")

    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return '\n'.join(full_text)

def routing(file_name, folder, database_name):
    path = "../../data_base/"+ database_name + "/" + folder + "/" + file_name
    file_ext = file_name[file_name.rindex('.') + 1:].lower()
    
    #text files kinda
    if file_ext == "pdf":
        return pdf_to_text(path)

    elif file_ext == "txt":
        return txt_to_text(path)

    elif file_ext == "docx":
        return docx_to_text(path)

    elif file_ext == "doc":
        return doc_to_text(path)
    
    elif file_ext == "csv":
        return csv_to_text(path)
    
    elif file_ext == "pptx":
        return pptx_to_text(path)
    
    elif file_ext == "ppt":
        return ppt_to_text(path)

    genai_id = upload(file_name, folder, database_name)

    #exle files 
    if file_ext == "xlsx":
        return ai_assist(genai_id)
    elif file_ext == "xls":
        return ai_assist(genai_id)
    elif file_ext == "ods":
        return ai_assist(genai_id)

    #image files
    elif file_ext == "jpg":
        return ai_assist(genai_id)
    elif file_ext == "png":
        return ai_assist(genai_id)
    elif file_ext == "jpeg":
        return ai_assist(genai_id)
    elif file_ext == "gif":
        return ai_assist(genai_id)

    #video files
    elif file_ext == "mp4":
        return ai_assist(genai_id)
    elif file_ext == "avi":
        return ai_assist(genai_id)
    
    cleanup(genai_id)
    